"""
Document engine — VerifyDoc.

Extracts text AND embedded images from PDF / DOCX / PPTX files, then runs:
  - image_engine.analyze_image_from_bytes on every extracted image
    (full pipeline, including the L11-L14 physical-consistency layers —
    PAFRA / BDIS / SSWDP / QESM — the same "advanced physics" checks used
    on standalone image scans, no separate/lite path).
  - text_engine.analyze_text on the extracted text.
  - plagiarism_engine.analyze_plagiarism_risk on the extracted text.

Image detection and text detection run IN PARALLEL via a shared
ThreadPoolExecutor (both engines are CPU-bound and release the GIL during
their numpy/torch inner loops, so real wall-clock parallelism is achieved
even though this module itself is synchronous underneath the async endpoint).

If the document contains no extractable images, only the text branch runs
(and vice versa for image-only PDFs with no OCR-able text) — the tool
never blocks waiting on a branch that has nothing to do.
"""

from __future__ import annotations

import io
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Dict, List, Tuple

MAX_IMAGES_PER_DOCUMENT = 20   # cap — avoid pathological documents stalling a request
MAX_TEXT_CHARS = 20000         # matches MAX_TEXT_LENGTH ballpark used elsewhere
MIN_IMAGE_BYTES = 3 * 1024     # skip tiny embedded icons/bullets/logos (noise, not content)


class UnsupportedDocumentError(Exception):
    pass


# ── Extraction ──────────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> Tuple[str, List[bytes], int]:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts: List[str] = []
    images: List[bytes] = []

    for page in doc:
        text_parts.append(page.get_text())
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base = doc.extract_image(xref)
                img_bytes = base.get("image")
                if img_bytes and len(img_bytes) >= MIN_IMAGE_BYTES:
                    images.append(img_bytes)
            except Exception:
                continue
            if len(images) >= MAX_IMAGES_PER_DOCUMENT:
                break
        if len(images) >= MAX_IMAGES_PER_DOCUMENT:
            break

    page_count = doc.page_count
    doc.close()
    return "\n\n".join(text_parts), images[:MAX_IMAGES_PER_DOCUMENT], page_count


def _extract_docx(file_bytes: bytes) -> Tuple[str, List[bytes], int]:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(file_bytes))
    text_parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text_parts.append(cell.text)

    images: List[bytes] = []
    for rel in document.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                if len(img_bytes) >= MIN_IMAGE_BYTES:
                    images.append(img_bytes)
            except Exception:
                continue
            if len(images) >= MAX_IMAGES_PER_DOCUMENT:
                break

    return "\n\n".join(text_parts), images[:MAX_IMAGES_PER_DOCUMENT], len(document.paragraphs)


def _extract_pptx(file_bytes: bytes) -> Tuple[str, List[bytes], int]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts: List[str] = []
    images: List[bytes] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        text_parts.append(line)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) >= MIN_IMAGE_BYTES:
                        images.append(img_bytes)
                except Exception:
                    continue
                if len(images) >= MAX_IMAGES_PER_DOCUMENT:
                    break

    slide_count = len(prs.slides)
    return "\n\n".join(text_parts), images[:MAX_IMAGES_PER_DOCUMENT], slide_count


def extract_document(file_bytes: bytes, content_type: str, filename: str) -> Dict[str, Any]:
    """Dispatch extraction by content-type/extension. Returns text, images, page/slide count, doc type."""
    name = (filename or "").lower()
    ct = (content_type or "").lower()

    if "pdf" in ct or name.endswith(".pdf"):
        text, images, units = _extract_pdf(file_bytes)
        doc_type = "pdf"
    elif "wordprocessingml" in ct or name.endswith(".docx"):
        text, images, units = _extract_docx(file_bytes)
        doc_type = "docx"
    elif "presentationml" in ct or name.endswith(".pptx"):
        text, images, units = _extract_pptx(file_bytes)
        doc_type = "pptx"
    else:
        raise UnsupportedDocumentError(
            f"Unsupported document type '{content_type}' / '{filename}'. "
            f"VerifyDoc currently supports PDF, DOCX, and PPTX."
        )

    return {
        "text": text.strip(),
        "images": images,
        "units": units,          # pages (pdf/docx-ish) or slides (pptx)
        "doc_type": doc_type,
        "image_count": len(images),
    }


# ── Parallel detection orchestration ───────────────────────────────────────

def analyze_document_from_bytes(
    file_bytes: bytes,
    content_type: str,
    filename: str,
    job_id: str = "",
) -> Dict[str, Any]:
    from engines.image_engine import analyze_image_from_bytes
    from engines.text_engine import analyze_text
    from engines.plagiarism_engine import analyze_plagiarism_risk, document_fingerprint

    t0 = time.time()

    extraction = extract_document(file_bytes, content_type, filename)
    text = extraction["text"][:MAX_TEXT_CHARS]
    images = extraction["images"]

    text_result: Dict[str, Any] | None = None
    plagiarism_result: Dict[str, Any] | None = None
    image_results: List[Dict[str, Any]] = []

    # Both branches (text detection+plagiarism, and N image scans) are
    # dispatched into the SAME pool so they genuinely run concurrently —
    # "find images first, and if not there then text detection works" per
    # product spec actually means: try both, use whichever branch has
    # content, and don't make one wait on the other.
    with ThreadPoolExecutor(max_workers=max(4, len(images) + 2)) as pool:
        futures = {}

        if text and len(text) >= 30:
            futures[pool.submit(analyze_text, text=text, job_id=job_id)] = ("text", -1)
            futures[pool.submit(analyze_plagiarism_risk, text)] = ("plagiarism", -1)

        for idx, img_bytes in enumerate(images):
            futures[pool.submit(
                analyze_image_from_bytes, img_bytes, "image/png", f"{job_id}_img{idx}"
            )] = ("image", idx)

        for future in as_completed(futures):
            kind, idx = futures[future]
            try:
                result = future.result()
            except Exception as e:
                result = {"status": "error", "error": str(e)}

            if kind == "text":
                text_result = result
            elif kind == "plagiarism":
                plagiarism_result = result
            elif kind == "image":
                image_results.append({"index": idx, **result})

    image_results.sort(key=lambda r: r["index"])

    # ── Composite verdict ───────────────────────────────────────────────
    ai_image_count = sum(
        1 for r in image_results
        if r.get("verdict") == "AI" or r.get("classification") == "AI"
    )
    text_is_ai = bool(text_result and text_result.get("verdict") == "AI")
    plag_high = bool(plagiarism_result and plagiarism_result.get("risk_level") == "HIGH")

    flags = []
    if ai_image_count > 0:
        flags.append(f"{ai_image_count} of {len(image_results)} embedded image(s) flagged as AI-generated")
    if text_is_ai:
        flags.append("document text flagged as likely AI-generated")
    if plag_high:
        flags.append("high originality-risk indicators in document text")

    if flags:
        composite_verdict = "FLAGGED"
        composite_summary = "; ".join(flags) + "."
    elif not images and not text:
        composite_verdict = "NO_CONTENT"
        composite_summary = "No extractable text or images were found in this document."
    else:
        composite_verdict = "CLEAN"
        composite_summary = "No AI-generation or originality-risk indicators found."

    return {
        "status": "ok",
        "document_type": extraction["doc_type"],
        "units_analyzed": extraction["units"],
        "document_fingerprint": document_fingerprint(text) if text else None,
        "has_text": bool(text and len(text) >= 30),
        "has_images": len(images) > 0,
        "image_count": len(images),
        "text_analysis": text_result,
        "image_analyses": image_results,
        "plagiarism_analysis": plagiarism_result,
        "composite_verdict": composite_verdict,
        "composite_summary": composite_summary,
        "processing_time_ms": round((time.time() - t0) * 1000, 1),
    }
